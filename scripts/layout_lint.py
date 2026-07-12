#!/usr/bin/env python3
"""Mechanical layout lint for generated .pptx files.

Statically analyzes a .pptx and flags likely layout problems before visual QA:

  OUT-OF-SLIDE      shape extends beyond the slide canvas            (ERROR)
  TEXT-OVERFLOW     estimated text height exceeds its box            (WARN/ERROR)
  NOWRAP-WIDTH      wrap disabled and the line is wider than the box (WARN)
  TEXT-OVERLAP      two text boxes intersect                         (WARN/ERROR)
  EDGE-MARGIN       text closer than 0.3" to a slide edge            (WARN)
  TINY-FONT         < 10pt body text / < 8pt table text              (ERROR)
  TABLE-COLW        sum of column widths != table width (CJK garble) (ERROR)
  MISSING-CJK-FONT  CJK run without an explicit fontFace             (WARN, ERROR in tables)
  EMOJI             emoji characters (render as boxes in LibreOffice)(ERROR)
  SPARSE            slide with very few items and little text        (WARN)
  LINE-START-KINSOKU 行頭禁則文字(。、ー等)が折り返しで行頭に来る推定   (WARN)

Text fit is ESTIMATED from character widths (CJK = fontSize pt wide, Latin
half that), so expect some false positives — treat WARN as "look at this
slide in the rendered image", not as a verdict. ERROR findings are almost
always real and must be fixed.

Usage:
    python layout_lint.py deck.pptx [more.pptx ...] [--json]

Exit code: 1 if any ERROR was found, else 0.
"""

import argparse
import json
import math
import sys
import unicodedata

# Windows consoles default to cp932/cp1252 and choke on findings text
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Emu

EMU_PER_IN = 914400

# Estimation constants (must match the Size Math rules in SKILL.md)
DEFAULT_FONT_PT = 18.0          # pptxgenjs default when fontSize is omitted
DEFAULT_LINE_SPACING = 1.2
OVERFLOW_WARN_RATIO = 1.10
OVERFLOW_ERROR_RATIO = 1.40
OVERLAP_WARN_RATIO = 0.12       # intersection / smaller-box area
OVERLAP_ERROR_RATIO = 0.40
EDGE_MARGIN_IN = 0.20
MIN_BODY_PT = 10.0
MIN_TABLE_PT = 8.0
SPARSE_MAX_ITEMS = 3            # flag if items <= this AND chars below threshold
SPARSE_MAX_CHARS = 120
# Default text-frame insets (EMU) when not set explicitly
DEF_INSET_LR = 91440
DEF_INSET_TB = 45720

# 行頭禁則文字（行頭に来てはいけない: 句読点・閉じ括弧・長音・小書き仮名・中点）。
# 半角ピリオド/カンマはファイルパスや英文で多用されるため除外（全角のみ対象）。
LINE_START_FORBIDDEN = set(
    "。、，．・：；！？”’）」』】］｝〉》〕…ー〜～"
    "ぁぃぅぇぉっゃゅょゎァィゥェォッャュョヮ"
)


def _in(emu):
    """EMU -> inches (float)."""
    return emu / EMU_PER_IN if emu is not None else None


def eff_char_width(ch):
    """Width of a char relative to fontSize-pt em: CJK fullwidth = 1, else 0.5."""
    if unicodedata.east_asian_width(ch) in ("W", "F"):
        return 1.0
    return 0.5


def has_cjk(text):
    return any(unicodedata.east_asian_width(c) in ("W", "F") for c in text)


def has_emoji(text):
    return any(ord(c) >= 0x1F000 for c in text)


def wrap_lines(sub, chars_per_line):
    """Greedy character wrap ignoring kinsoku. Returns display-line strings.

    Mirrors the renderer's CJK any-character wrapping without applying the
    line-start-forbidden push/hang, so a forbidden char appearing at the head
    of a wrapped line here is exactly what the renderer must fix up."""
    lines, cur, cur_w = [], "", 0.0
    for ch in sub:
        cw = eff_char_width(ch)
        if cur and cur_w + cw > chars_per_line + 1e-9:
            lines.append(cur)
            cur, cur_w = ch, cw
        else:
            cur += ch
            cur_w += cw
    if cur:
        lines.append(cur)
    return lines if lines else [""]


def run_font_size_pt(run, para):
    if run.font.size is not None:
        return run.font.size.pt
    if para.font.size is not None:
        return para.font.size.pt
    return None


def run_font_name(run, para):
    return run.font.name or para.font.name


def para_line_height_in(para, font_pt):
    ls = para.line_spacing
    if ls is None:
        return font_pt * DEFAULT_LINE_SPACING / 72.0
    if isinstance(ls, float):
        return font_pt * ls / 72.0
    return ls.pt / 72.0  # exact spacing as Length


def estimate_text_height_in(tf, usable_w_in):
    """Estimated height (inches) the text needs inside `usable_w_in`.

    Returns (height, detail_string) or (None, reason) when not estimable.
    """
    if usable_w_in is None or usable_w_in <= 0.05:
        return None, "no usable width"
    total = 0.0
    lines_total = 0
    max_pt = 0.0
    for para in tf.paragraphs:
        sizes = [run_font_size_pt(r, para) for r in para.runs]
        sizes = [s for s in sizes if s]
        font_pt = max(sizes) if sizes else (
            para.font.size.pt if para.font.size else DEFAULT_FONT_PT)
        max_pt = max(max_pt, font_pt)
        chars_per_line = usable_w_in * 72.0 / font_pt
        # '\v' marks explicit <a:br> line breaks in paragraph text
        sublines = para.text.split("\v") if para.text else [""]
        n_lines = 0
        for sub in sublines:
            w = sum(eff_char_width(c) for c in sub)
            n_lines += max(1, math.ceil(w / chars_per_line)) if w else 1
        lines_total += n_lines
        total += n_lines * para_line_height_in(para, font_pt)
    return total, f"~{lines_total} line(s) @{max_pt:.0f}pt"


def shape_label(shape):
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text.replace("\n", " ").replace("\v", " ")
    if len(text) > 24:
        text = text[:24] + "…"
    return f'"{text}"' if text else f"<{shape.shape_type}>"


class Linter:
    def __init__(self, path):
        self.path = path
        self.findings = []  # dicts: slide, level, code, msg

    def add(self, slide_no, level, code, msg):
        finding = {"slide": slide_no, "level": level, "code": code, "msg": msg}
        if finding not in self.findings:  # same issue across runs/cells
            self.findings.append(finding)

    def lint(self):
        prs = Presentation(self.path)
        sw, sh = _in(prs.slide_width), _in(prs.slide_height)
        for idx, slide in enumerate(prs.slides, start=1):
            self.lint_slide(idx, slide, sw, sh)
        return self.findings

    def lint_slide(self, no, slide, sw, sh):
        text_rects = []   # (l, t, r, b, label) for overlap pass
        n_items = 0
        n_chars = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.add(no, "WARN", "GROUP",
                         "group shape skipped (positions not checked)")
                continue
            # スライド番号などの自動フィールドプレースホルダは内容・フォントがレンダラ任せ
            # （マスター由来で python-pptx からは読めない）。静的な fit/font 判定の対象外。
            if shape.is_placeholder and shape.placeholder_format is not None:
                try:
                    if shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                        continue
                except (KeyError, ValueError):
                    pass
            l, t = _in(shape.left), _in(shape.top)
            w, h = _in(shape.width), _in(shape.height)
            label = shape_label(shape)

            if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART):
                n_items += 1

            # bounds — text off-canvas is a real bug; decorative shape bleed
            # (circles/bars running off the edge) is a legitimate design choice
            if None not in (l, t, w, h):
                if l < -0.01 or t < -0.01 or l + w > sw + 0.01 or t + h > sh + 0.01:
                    has_text = shape.has_text_frame and shape.text_frame.text.strip()
                    self.add(no, "ERROR" if has_text else "WARN", "OUT-OF-SLIDE",
                             f"{label} at ({l:.2f},{t:.2f}) {w:.2f}x{h:.2f}in "
                             f"exceeds {sw:.2f}x{sh:.2f}in canvas"
                             + ("" if has_text else " (fine if decorative bleed)"))

            if shape.has_table:
                n_items += 1
                self.lint_table(no, shape, label)
                continue
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            text = tf.text
            if text.strip():
                n_items += 1
                n_chars += len(text.replace("\n", "").replace("\v", ""))
                if None not in (l, t, w, h):
                    text_rects.append((l, t, l + w, t + h, label))
                    # top edge exempt: header/title bands legitimately sit high
                    if (l < EDGE_MARGIN_IN
                            or sw - (l + w) < EDGE_MARGIN_IN
                            or sh - (t + h) < EDGE_MARGIN_IN):
                        self.add(no, "WARN", "EDGE-MARGIN",
                                 f"{label} is within {EDGE_MARGIN_IN}\" of a "
                                 "left/right/bottom slide edge")
            self.lint_text_runs(no, tf, label, in_table=False)
            if text.strip() and None not in (w, h):
                self.lint_text_fit(no, tf, w, h, label)
                self.lint_kinsoku(no, tf, w, h, label)

        self.lint_overlaps(no, text_rects)
        if n_items <= SPARSE_MAX_ITEMS and n_chars < SPARSE_MAX_CHARS:
            self.add(no, "WARN", "SPARSE",
                     f"only {n_items} item(s), {n_chars} chars — fine for "
                     "hero/section-divider/quote, otherwise add substance")

    def lint_text_fit(self, no, tf, w, h, label):
        # autofit (shrink/resize) makes static estimation meaningless
        if tf.auto_size is not None:
            return
        ml = _in(tf.margin_left) if tf.margin_left is not None else _in(DEF_INSET_LR)
        mr = _in(tf.margin_right) if tf.margin_right is not None else _in(DEF_INSET_LR)
        mt = _in(tf.margin_top) if tf.margin_top is not None else _in(DEF_INSET_TB)
        mb = _in(tf.margin_bottom) if tf.margin_bottom is not None else _in(DEF_INSET_TB)
        usable_w = w - ml - mr
        usable_h = h - mt - mb
        if tf.word_wrap is False:
            # single-line width check instead of wrap estimation
            for para in tf.paragraphs:
                sizes = [run_font_size_pt(r, para) for r in para.runs]
                sizes = [s for s in sizes if s]
                font_pt = max(sizes) if sizes else DEFAULT_FONT_PT
                for sub in (para.text.split("\v") if para.text else []):
                    line_w = sum(eff_char_width(c) for c in sub) * font_pt / 72.0
                    if line_w > usable_w * 1.05:
                        self.add(no, "WARN", "NOWRAP-WIDTH",
                                 f"{label}: line needs ~{line_w:.2f}\" but box "
                                 f"offers {usable_w:.2f}\" (wrap disabled)")
            return
        needed, detail = estimate_text_height_in(tf, usable_w)
        if needed is None or usable_h <= 0:
            return
        ratio = needed / usable_h
        if ratio > OVERFLOW_ERROR_RATIO:
            level = "ERROR"
        elif ratio > OVERFLOW_WARN_RATIO:
            level = "WARN"
        else:
            return
        self.add(no, level, "TEXT-OVERFLOW",
                 f"{label}: {detail} needs ~{needed:.2f}\" but box offers "
                 f"{usable_h:.2f}\" ({ratio:.0%})")

    def lint_kinsoku(self, no, tf, w, h, label):
        # 折り返しで行頭禁則文字（。、ー 等）が行頭に来る箇所を推定検出。実レンダラは
        # 禁則処理で前行へ追い込む/ぶら下げるため、ここで挙がる箇所は折り返しが乱れて
        # 見えやすい。明示改行や字句調整で句読点を前行末に固定して潰す。
        if tf.auto_size is not None or tf.word_wrap is False:
            return
        ml = _in(tf.margin_left) if tf.margin_left is not None else _in(DEF_INSET_LR)
        mr = _in(tf.margin_right) if tf.margin_right is not None else _in(DEF_INSET_LR)
        usable_w = w - ml - mr
        if usable_w <= 0.05:
            return
        for para in tf.paragraphs:
            sizes = [run_font_size_pt(r, para) for r in para.runs]
            sizes = [s for s in sizes if s]
            font_pt = max(sizes) if sizes else (
                para.font.size.pt if para.font.size else DEFAULT_FONT_PT)
            chars_per_line = usable_w * 72.0 / font_pt
            if chars_per_line < 2:
                continue
            for sub in (para.text.split("\v") if para.text else []):
                lines = wrap_lines(sub, chars_per_line)
                for li in range(1, len(lines)):
                    if lines[li] and lines[li][0] in LINE_START_FORBIDDEN:
                        prev_tail = lines[li - 1][-8:]
                        self.add(no, "WARN", "LINE-START-KINSOKU",
                                 f'{label}: 折り返しで行頭禁則 "{lines[li][0]}" が行頭に'
                                 f'来る推定（"…{prev_tail}" の直後）')
                        break  # one report per paragraph is enough

    def lint_text_runs(self, no, tf, label, in_table):
        min_pt = MIN_TABLE_PT if in_table else MIN_BODY_PT
        for para in tf.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                if has_emoji(run.text):
                    self.add(no, "ERROR", "EMOJI",
                             f"{label}: emoji in \"{run.text[:20]}\" — "
                             "LibreOffice renders emoji as boxes")
                size = run_font_size_pt(run, para)
                if size is not None and size < min_pt:
                    self.add(no, "ERROR", "TINY-FONT",
                             f"{label}: {size:.1f}pt text (minimum "
                             f"{min_pt:.0f}pt{' in tables' if in_table else ''})")
                if has_cjk(run.text) and run_font_name(run, para) is None:
                    self.add(no, "ERROR" if in_table else "WARN",
                             "MISSING-CJK-FONT",
                             f"{label}: CJK text without explicit fontFace"
                             f"{' (garbles in PowerPoint tables)' if in_table else ''}")

    def lint_table(self, no, shape, label):
        table = shape.table
        col_sum = sum(c.width for c in table.columns)
        if abs(col_sum - shape.width) > Emu(int(0.01 * EMU_PER_IN)):
            self.add(no, "ERROR", "TABLE-COLW",
                     f"table: colW sums to {_in(col_sum):.2f}\" but table w is "
                     f"{_in(shape.width):.2f}\" — PowerPoint recalculates and "
                     "may garble CJK")
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame.text.strip():
                    self.lint_text_runs(no, cell.text_frame,
                                        "table cell", in_table=True)

    def lint_overlaps(self, no, rects):
        for i in range(len(rects)):
            for j in range(i + 1, len(rects)):
                l1, t1, r1, b1, a = rects[i]
                l2, t2, r2, b2, b = rects[j]
                iw = min(r1, r2) - max(l1, l2)
                ih = min(b1, b2) - max(t1, t2)
                if iw <= 0 or ih <= 0:
                    continue
                inter = iw * ih
                smaller = min((r1 - l1) * (b1 - t1), (r2 - l2) * (b2 - t2))
                if smaller <= 0:
                    continue
                ratio = inter / smaller
                if ratio > OVERLAP_ERROR_RATIO:
                    level = "ERROR"
                elif ratio > OVERLAP_WARN_RATIO:
                    level = "WARN"
                else:
                    continue
                self.add(no, level, "TEXT-OVERLAP",
                         f"{a} and {b} overlap by {ratio:.0%} of the smaller box")


def main():
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("files", nargs="+", help=".pptx file(s) to lint")
    ap.add_argument("--json", action="store_true", help="machine-readable output")
    args = ap.parse_args()

    any_error = False
    all_results = {}
    for path in args.files:
        findings = Linter(path).lint()
        all_results[path] = findings
        errors = sum(1 for f in findings if f["level"] == "ERROR")
        warns = sum(1 for f in findings if f["level"] == "WARN")
        any_error |= errors > 0
        if args.json:
            continue
        print(f"== {path} ==")
        last_slide = None
        for f in sorted(findings, key=lambda f: f["slide"]):
            if f["slide"] != last_slide:
                print(f" slide {f['slide']}:")
                last_slide = f["slide"]
            print(f"   [{f['level']}] {f['code']}: {f['msg']}")
        print(f" {errors} error(s), {warns} warning(s)")
        if warns:
            print("  (WARN = estimation heuristic — verify against the rendered "
                  "image before changing code)")
    if args.json:
        print(json.dumps(all_results, ensure_ascii=False, indent=1))
    sys.exit(1 if any_error else 0)


if __name__ == "__main__":
    main()
