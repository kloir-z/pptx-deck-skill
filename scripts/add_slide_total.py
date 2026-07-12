#!/usr/bin/env python3
"""Append a fixed " / <total>" to every PowerPoint slide-number field.

PowerPoint has no total-slides field (Word does), so the current number must stay
an auto-updating `slidenum` field while the total is embedded as a literal run in
the **same paragraph** — a separate text box would not share the field's vertical
anchor and would misalign vertically in PowerPoint.

This opens the saved .pptx, finds each SLIDE_NUMBER placeholder, and inserts an
`<a:r>` (inheriting the field's run properties) right after the `<a:fld>` in its
paragraph. Idempotent: re-running does not duplicate the total.

Run this once after generating the .pptx (e.g. from the generator's writeFile
callback, or as a build step). Make the slide-number box wide enough that
`n / total` fits on one line.

Usage: python add_slide_total.py deck.pptx
"""
import copy
import sys

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn


def add_slide_totals(path):
    prs = Presentation(path)
    total = len(prs.slides)
    added = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not (sh.is_placeholder and sh.placeholder_format is not None):
                continue
            try:
                if sh.placeholder_format.type != PP_PLACEHOLDER.SLIDE_NUMBER:
                    continue
            except (KeyError, ValueError):
                continue
            p = sh.text_frame._txBody.find(qn("a:p"))
            if p is None:
                continue
            fld = p.find(qn("a:fld"))
            if fld is None:
                continue
            # 冪等: 既に固定テキスト（「/」を含む run）があればスキップ
            if any("/" in (t.text or "") for t in p.iter(qn("a:t"))):
                continue
            # fld の rPr を引き継いだ run を作り、fld の直後（endParaRPr の手前）に並べる
            r = p.makeelement(qn("a:r"), {})
            rpr = fld.find(qn("a:rPr"))
            if rpr is not None:
                r.append(copy.deepcopy(rpr))
            t = p.makeelement(qn("a:t"), {})
            t.text = f" / {total}"
            r.append(t)
            endpr = p.find(qn("a:endParaRPr"))
            if endpr is not None:
                endpr.addprevious(r)
            else:
                p.append(r)
            added += 1
    prs.save(path)
    print(f'add_slide_total: appended " / {total}" to {added} slide-number field(s)')


if __name__ == "__main__":
    add_slide_totals(sys.argv[1])
